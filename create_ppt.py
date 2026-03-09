"""
Generate presentation: Worm Robot Pipe Crawl — Progress Report
Expanded version with V1→V2→V3 evolution & detailed Riddle 2025 analysis.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ─── Paths ───────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
IMG_V3 = os.path.join(BASE, "record", "v3", "videos")
IMG_RIDDLE = r"D:\inovxio\paper\myPHS\Riddle_2025\images"
OUT = os.path.join(BASE, "record", "worm_pipe_crawl_report.pptx")

# ─── Colors ──────────────────────────────────────────────────────────
DARK_BG   = RGBColor(0x1a, 0x1a, 0x2e)
BLUE_ACC  = RGBColor(0x00, 0x96, 0xc7)
GREEN_ACC = RGBColor(0x2d, 0xce, 0x89)
RED_ACC   = RGBColor(0xf5, 0x36, 0x5c)
ORANGE_ACC= RGBColor(0xfb, 0x8c, 0x00)
PURPLE_ACC= RGBColor(0xa2, 0x6d, 0xff)
WHITE     = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY= RGBColor(0xcc, 0xcc, 0xcc)
MID_GRAY  = RGBColor(0x88, 0x88, 0x88)
DARK_GRAY = RGBColor(0x33, 0x33, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18,
             color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=WHITE, line_spacing=1.3, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (txt, sz, clr, bld) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz if sz else font_size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bld if bld else False
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 0.6)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, bullet_color=BLUE_ACC, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            txt, clr = item
        else:
            txt, clr = item, color
        p.text = txt
        p.font.size = Pt(font_size)
        p.font.color.rgb = clr
        p.font.name = font_name
        p.space_after = Pt(4)
        pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        for child in list(pPr):
            if child.tag.endswith('buChar') or child.tag.endswith('buNone'):
                pPr.remove(child)
        pPr.append(buChar)
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': str(bullet_color)})
        buClr.append(srgb)
        for child in list(pPr):
            if child.tag.endswith('buClr'):
                pPr.remove(child)
        pPr.append(buClr)
    return txBox


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_image_safe(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kwargs = {}
        if width: kwargs['width'] = Inches(width)
        if height: kwargs['height'] = Inches(height)
        return slide.shapes.add_picture(path, Inches(left), Inches(top), **kwargs)
    return None


def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=BLUE_ACC, cell_color=DARK_GRAY,
              text_color=WHITE, font_size=12):
    actual_rows = len(data)
    actual_cols = len(data[0]) if data else cols
    ts = slide.shapes.add_table(actual_rows, actual_cols, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    table = ts.table
    for r in range(actual_rows):
        for c in range(actual_cols):
            cell = table.cell(r, c)
            cell.text = str(data[r][c]) if c < len(data[r]) else ""
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = text_color
                p.font.name = "Microsoft YaHei"
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color if r == 0 else cell_color
    return ts


# =====================================================================
# SLIDE 1: Title
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_rect(slide, 0, 7.42, 13.333, 0.08, BLUE_ACC)

add_text(slide, 1.5, 1.5, 10, 1.2,
         "\u8718\u86f7\u673a\u5668\u4eba\u7ba1\u9053\u722c\u884c\u4eff\u771f",
         font_size=42, color=WHITE, bold=True)
add_text(slide, 1.5, 2.5, 10, 0.8,
         "Worm Robot Pipe Crawl Simulation",
         font_size=28, color=BLUE_ACC)
add_text(slide, 1.5, 3.3, 10, 0.8,
         "Passive Turning via Wall Constraints in MuJoCo",
         font_size=22, color=LIGHT_GRAY)

add_multiline(slide, 1.5, 4.3, 10, 2.5, [
    ("5\u6bb5\u53d8\u6001\u8718\u86f7  |  MuJoCo\u7269\u7406\u5f15\u64ce  |  \u5f39\u7c27\u94a2\u7efc\u7ebf\u6a21\u578b", 16, LIGHT_GRAY, False),
    ("", 8, None, False),
    ("\u5173\u952e\u7ed3\u679c: \u7403\u94f0\u7ea6\u675f <connect> \u5b9e\u73b0 88.7\u00b0 \u88ab\u52a8\u8f6c\u5f2f", 20, GREEN_ACC, True),
    ("", 8, None, False),
    ("\u53c2\u8003\u6587\u732e: Riddle et al. 2025, Bioinspir. Biomim.", 14, MID_GRAY, False),
    ("             Zhan et al. 2019, IJRR", 14, MID_GRAY, False),
    ("", 8, None, False),
    ("2026-03-01", 14, MID_GRAY, False),
])


# =====================================================================
# SLIDE 2: Outline
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 5, 0.6, "\u5185\u5bb9\u5927\u7eb2  Outline", font_size=32, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 2, 0.04, BLUE_ACC)

items = [
    ("\u2776  \u80cc\u666f\u4e0e\u52a8\u673a  Background & Motivation", WHITE),
    ("\u2777  \u8bba\u6587\u5206\u6790: Zhan 2019 \u6b65\u6001\u7406\u8bba", WHITE),
    ("\u2778  \u8bba\u6587\u5206\u6790: Riddle 2025 \u7ba1\u9053\u722c\u884c (\u8be6\u7ec6)", BLUE_ACC),
    ("\u2779  \u6a21\u578b\u6f14\u8fdb: V1 \u2192 V2 \u2192 V3", PURPLE_ACC),
    ("\u277a  V3\u6a21\u578b\u67b6\u6784\u4e0e\u5355\u6bb5\u9a8c\u8bc1", WHITE),
    ("\u277b  \u5b9e\u9a8c\u7ed3\u679c\u603b\u89c8", WHITE),
    ("\u277c  \u7a81\u7834: <connect> vs <weld> \u7ea6\u675f", GREEN_ACC),
    ("\u277d  \u4e0b\u4e00\u6b65\u8ba1\u5212  Next Steps", ORANGE_ACC),
]
for i, (item, clr) in enumerate(items):
    add_text(slide, 1.5, 1.4 + i * 0.7, 10, 0.6, item,
             font_size=20, color=clr, bold=(clr == GREEN_ACC))


# =====================================================================
# SLIDE 3: Background
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 8, 0.6, "\u80cc\u666f\u4e0e\u52a8\u673a  Background", font_size=32, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 3.5, 0.04, BLUE_ACC)

add_bullet_list(slide, 0.8, 1.4, 6, 5.5, [
    "\u8718\u86f7\u4eff\u751f\u673a\u5668\u4eba\u9002\u5408\u53d7\u9650\u7a7a\u95f4 (\u7ba1\u9053/\u96a7\u9053/\u5e9f\u589f)",
    "\u8815\u52a8\u8fd0\u52a8 (Peristaltic): \u9006\u884c\u808c\u8089\u6ce2 \u2192 \u524d\u8fdb\u8fd0\u52a8",
    "\u88ab\u52a8\u67d4\u987a\u6027\u5141\u8bb8\u65e0\u4e3b\u52a8\u8f6c\u5411\u7684\u7ba1\u9053\u5bfc\u822a",
    "\u6311\u6218: \u51c6\u786e\u5efa\u6a21\u8f6f\u4f53+\u63a5\u89e6\u52a8\u529b\u5b66",
    "\u76ee\u6807: MuJoCo 5\u6bb5\u8718\u86f7\u5728\u5f2f\u7ba1\u4e2d\u7684\u88ab\u52a8\u8f6c\u5f2f\u4eff\u771f",
], font_size=17)

add_rect(slide, 7.3, 1.4, 5.5, 2.5, DARK_GRAY)
add_text(slide, 7.5, 1.5, 5, 0.5, "\u6838\u5fc3\u95ee\u9898  Core Questions", font_size=20, color=ORANGE_ACC, bold=True)
add_bullet_list(slide, 7.5, 2.0, 5, 2, [
    "\u8718\u86f7\u80fd\u5426\u5728\u5f2f\u7ba1\u4e2d\u88ab\u52a8\u8f6c\u5f2f?",
    "\u8f6c\u5f2f\u80fd\u529b\u53d6\u51b3\u4e8e\u4ec0\u4e48: \u521a\u5ea6? \u5173\u8282\u7c7b\u578b?",
    "\u5982\u4f55\u5728MuJoCo\u4e2d\u5efa\u6a21\u5f39\u7c27\u94a2\u7efc\u7ebf?",
], font_size=15, bullet_color=ORANGE_ACC)

add_image_safe(slide, os.path.join(IMG_RIDDLE,
    "59e16b0384219388882768541b392048cfe3bc07d2e9ef810b1d3938e959ee1f.jpg"),
    7.3, 4.2, width=5.5)


# =====================================================================
# SLIDE 4: Zhan 2019 Gait Theory
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u8bba\u6587: Zhan et al. 2019 \u2014 \u6b65\u6001\u7406\u8bba (IJRR)",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 5, 0.04, BLUE_ACC)

add_text(slide, 0.8, 1.3, 5, 0.5, "4\u79cd\u6bb5\u72b6\u6001  4 Segment States",
         font_size=20, color=ORANGE_ACC, bold=True)
state_data = [
    ["\u72b6\u6001", "\u7f16\u7801", "\u52a8\u4f5c", "\u89d2\u5ea6"],
    ["\u653e\u677e Relaxed", "0", "\u73af\u808c\u6536\u7f29 (Ring ON)", "0\u00b0"],
    ["\u6536\u7f29 Contracted", "1", "\u5168\u8f74\u5411\u808c\u8089 ON", "0\u00b0"],
    ["\u5de6\u5f2f Left-bend", "2", "\u5de6\u4fa7\u8f74\u5411\u808c\u8089", "+\u03b8\u2080"],
    ["\u53f3\u5f2f Right-bend", "3", "\u53f3\u4fa7\u8f74\u5411\u808c\u8089", "-\u03b8\u2080"],
]
add_table(slide, 0.8, 1.8, 5.5, 2.0, 5, 4, state_data, font_size=13)

add_text(slide, 0.8, 4.2, 5, 0.5, "4\u79cd\u8fd0\u52a8\u6a21\u5f0f  4 Locomotion Modes",
         font_size=20, color=ORANGE_ACC, bold=True)
mode_data = [
    ["\u6a21\u5f0f", "\u6761\u4ef6", "\u8f68\u8ff9"],
    ["\u76f4\u7ebf Rectilinear", "n\u2082=n\u2083=0", "\u76f4\u7ebf"],
    ["\u4fa7\u7eed Sidewinding", "n\u2082=n\u2083\u22600", "\u5bf9\u89d2\u7ebf"],
    ["\u5706\u5f27 Circular", "n\u2082\u2260n\u2083, closed", "\u5f27\u7ebf"],
    ["\u6446\u7ebf Cycloid", "n\u2082\u2260n\u2083, open", "\u6446\u7ebf"],
]
add_table(slide, 0.8, 4.7, 5.5, 2.0, 5, 3, mode_data, font_size=13)

add_rect(slide, 7, 1.3, 5.8, 5.6, DARK_GRAY)
add_text(slide, 7.2, 1.4, 5.5, 0.5,
         "\u6b65\u6001\u53c2\u6570: {n\u2082, n\u2083, n\u2081 | nP}",
         font_size=20, color=GREEN_ACC, bold=True)
add_multiline(slide, 7.2, 2.0, 5.3, 4.5, [
    ("n\u2081 = \u951a\u5b9a\u6bb5 (State 1), n\u2082 = \u5de6\u5f2f\u6bb5, n\u2083 = \u53f3\u5f2f\u6bb5", 14, LIGHT_GRAY, False),
    ("nP = \u6ce2\u4f20\u64ad\u6b65\u957f,  TN = n\u2080+n\u2081+n\u2082+n\u2083+nP", 14, LIGHT_GRAY, False),
    ("", 8, None, False),
    ("\u6211\u4eec\u7684\u76f4\u7ebf\u6b65\u6001: {0,0,2|1}", 16, GREEN_ACC, True),
    ("  s\u2080 = [0,0,0,1,1] (\u5c3e\u2192\u5934)", 14, BLUE_ACC, False),
    ("  2-anchor, \u65e0\u5f2f\u66f2\u6bb5", 14, BLUE_ACC, False),
    ("", 8, None, False),
    ("\u5706\u5f27\u6b65\u6001: {1,0,1|1}", 16, GREEN_ACC, True),
    ("  s\u2080 = [2,0,0,0,1] (1\u4e2a\u5de6\u5f2f\u6bb5)", 14, BLUE_ACC, False),
    ("  \u03b8_T = \u03b8\u2080 = 18\u00b0/\u5468\u671f", 14, BLUE_ACC, False),
    ("", 8, None, False),
    ("5\u6bb5\u673a\u5668\u4eba: \u5171 25 \u79cd\u6709\u6548\u6b65\u6001", 16, ORANGE_ACC, True),
    ("  4 rectilinear + 3 sidewinding + 18 circular", 14, ORANGE_ACC, False),
])


# =====================================================================
# SLIDE 5: Riddle 2025 — Model Design (DETAILED)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 12, 0.6,
         "\u8bba\u6587: Riddle 2025 \u2014 \u6a21\u578b\u8bbe\u8ba1 (Model Design)",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 5, 0.04, BLUE_ACC)

# Left: mesh structure
add_text(slide, 0.8, 1.3, 5.5, 0.5, "\u7f51\u683c\u7ed3\u6784  Mesh Structure",
         font_size=18, color=ORANGE_ACC, bold=True)
add_bullet_list(slide, 0.8, 1.8, 5.5, 2.0, [
    "CMMWorm-SES: 12\u80a1\u87ba\u65cb\u805a\u4e59\u70ef\u7ba1 (6CW + 6CCW)",
    "\u7ba1\u6750: LDPE, \u5916\u5f84 0.25\", \u5185\u5f84 0.17\"",
    "MuJoCo\u7b49\u6548\u5b9e\u5fc3\u76f4\u5f84: 5.98mm (\u60ef\u6027\u77e9\u5339\u914d)",
    "\u9759\u6b62\u76f4\u5f84 320mm, \u6536\u7f29\u76f4\u5f84 210mm",
    "\u8d28\u91cf: 1.98 kg",
], font_size=14)

# Left: vertex joints
add_text(slide, 0.8, 3.5, 5.5, 0.5, "\u9876\u70b9\u5173\u8282  Vertex Joints",
         font_size=18, color=ORANGE_ACC, bold=True)
add_bullet_list(slide, 0.8, 4.0, 5.5, 2.0, [
    "\u7269\u7406\u673a\u5668\u4eba: \u5f00\u5c3e\u87ba\u4e1d\u9500\u5173\u8282 (eyelet screw)",
    ("MuJoCo: <connect> \u7ea6\u675f = \u7403\u94f0\u5173\u8282 (3DOF)", GREEN_ACC),
    ("solref = 0.001 (\u7d27\u675f\u7ea6, \u9632\u6b62\u6d6e\u52a8)", GREEN_ACC),
    "<exclude> \u6392\u9664\u4ea4\u53c9\u4f53\u7684\u63a5\u89e6\u529b",
    "\u7403\u5f62\u4f53: \u6807\u51c6 0.5cm/3.8g, \u4f20\u611f\u5668 2.24cm",
], font_size=14)

# Left: muscles
add_text(slide, 0.8, 5.7, 5.5, 0.5, "\u808c\u8089\u4e0e\u5f39\u7c27  Muscles & Springs",
         font_size=18, color=ORANGE_ACC, bold=True)
add_bullet_list(slide, 0.8, 6.1, 5.5, 1.2, [
    "\u6a61\u76ae\u7b4b\u5f39\u7c27: k=24.5 N/m, \u9759\u6b62\u957f 6.5cm",
    "Hill\u808c\u8089\u6a21\u578b: F\u2080=9N (0.3GPa) / 31N (2GPa)",
    "\u808c\u8089\u884c\u7a0b: 62-96cm (\u5bf9\u5e94\u5185\u63a5\u516d\u8fb9\u5f62\u51e0\u4f55)",
], font_size=14)

# Right: model images
add_image_safe(slide, os.path.join(IMG_RIDDLE,
    "f68939da3aeec56d6ccee422279247b5302ac3b339a26ebfbab50fe3dabf698b.jpg"),
    7, 1.2, width=3.0)
add_image_safe(slide, os.path.join(IMG_RIDDLE,
    "ac96ec9ed25d358c25f7ad73e4808dbfa1bfb17f81b21366a27c1583e744e17d.jpg"),
    10.2, 1.2, width=2.8)

# Right: material comparison
add_rect(slide, 7, 4.0, 5.8, 3.0, DARK_GRAY)
add_text(slide, 7.2, 4.1, 5.5, 0.5, "\u6750\u6599\u5bf9\u6bd4  Material Comparison",
         font_size=16, color=GREEN_ACC, bold=True)
mat_data = [
    ["", "Riddle (LDPE)", "\u6211\u4eec (Spring Steel)"],
    ["Young's Modulus", "0.3 GPa (tuned 2GPa)", "~200 GPa"],
    ["\u5173\u8282\u7c7b\u578b", "<connect> \u7403\u94f0", "<weld>\u2192<connect>"],
    ["\u7ed3\u6784", "12\u87ba\u65cb\u7f51\u683c", "8\u5e73\u884c\u7efc\u7ebf"],
    ["\u76f4\u5f84", "320 mm", "44 mm"],
    ["solref", "0.001", "0.005\u20130.01"],
]
add_table(slide, 7.1, 4.6, 5.5, 2.2, 6, 3, mat_data, font_size=11)


# =====================================================================
# SLIDE 6: Riddle 2025 — Results (DETAILED)
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 12, 0.6,
         "\u8bba\u6587: Riddle 2025 \u2014 \u5b9e\u9a8c\u7ed3\u679c (Results)",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 5, 0.04, BLUE_ACC)

# Speed vs Friction
add_text(slide, 0.8, 1.3, 5.5, 0.5,
         "\u2460 \u901f\u5ea6 vs \u6469\u64e6\u7cfb\u6570  Speed vs Friction",
         font_size=18, color=ORANGE_ACC, bold=True)
add_bullet_list(slide, 0.8, 1.8, 5.5, 1.8, [
    "\u6d4b\u8bd5\u8303\u56f4: \u03bc = 0.1 \u2013 0.9",
    ("2GPa (\u786c): \u9ad8\u6469\u64e6 = \u9ad8\u901f\u5ea6 (\u6bb5\u80fd\u79bb\u5730\u62ac\u8d77)", GREEN_ACC),
    ("0.3GPa (\u8f6f): \u9ad8\u6469\u64e6 = \u4f4e\u901f\u5ea6 (\u62d6\u5730\u8d70)", RED_ACC),
    "3\u00d71\u6ce2\u5f62\u59cb\u7ec8\u5feb\u4e8e 2\u00d71\u6ce2\u5f62",
    "\u5e73\u5730\u901f\u5ea6: 2.9\u20134.8 cm/s (2GPa, 3\u00d71)",
], font_size=13)

# Speed vs ROC — KEY FINDING
add_text(slide, 0.8, 3.5, 5.5, 0.5,
         "\u2461 \u901f\u5ea6 vs \u8f6c\u5f2f\u534a\u5f84  Speed vs ROC",
         font_size=18, color=ORANGE_ACC, bold=True)
add_rect(slide, 0.8, 4.0, 5.5, 2.5, RGBColor(0x1a, 0x2a, 0x1a))
add_multiline(slide, 1.0, 4.1, 5.2, 2.3, [
    ("\u4e34\u754c\u8f6c\u5f2f\u534a\u5f84 = 0.45m (ROC)", 17, GREEN_ACC, True),
    ("", 6, None, False),
    ("ROC > 0.45m: \u521a\u5ea6\u8d8a\u9ad8 \u2192 \u901f\u5ea6\u8d8a\u5feb", 14, LIGHT_GRAY, False),
    ("ROC \u2264 0.45m: \u521a\u5ea6\u5bf9\u901f\u5ea6\u65e0\u663e\u8457\u5f71\u54cd", 14, LIGHT_GRAY, False),
    ("", 6, None, False),
    ("\u6d4b\u8bd5\u8303\u56f4: 0.3m \u2013 0.9m ROC (\u03bc=0.3)", 13, LIGHT_GRAY, False),
    ("\u6700\u5c0f ROC 0.3m: \u4ec5 0.3GPa \u6a21\u578b\u80fd\u901a\u8fc7", 13, ORANGE_ACC, False),
    ("1GPa/2GPa \u6a21\u578b\u5728 ROC<0.35m \u65f6\u65e0\u6cd5\u5b8c\u6210\u8f6c\u5f2f", 13, RED_ACC, False),
])

# Joint Friction
add_text(slide, 0.8, 6.2, 5.5, 0.5,
         "\u2462 \u5173\u8282\u6469\u64e6\u662f\u5173\u952e\u56e0\u7d20",
         font_size=16, color=ORANGE_ACC, bold=True)
add_multiline(slide, 0.8, 6.55, 5.5, 0.8, [
    ("\u540d\u4e49 E=0.3GPa \u2192 \u7b49\u6548 E=2GPa (\u5173\u8282\u6469\u64e6\u589e\u52a0 3.3\u00d7 \u529b)", 12, LIGHT_GRAY, False),
    ("\u62e7\u7d27\u87ba\u4e1d\u540e\u529b\u4ece 16N \u589e\u5230 22.5N (+41%)", 12, LIGHT_GRAY, False),
])

# Right: ROC comparison images
add_image_safe(slide, os.path.join(IMG_RIDDLE,
    "233177cce3cd1b1905fb369c262123fbb9b63c4248a2caf39b383b24265ee964.jpg"),
    7, 1.2, width=5.8)

# Right: speed vs ROC graph
add_image_safe(slide, os.path.join(IMG_RIDDLE,
    "6ab38cf8b7e88ebc5b0e3e40ae28e9e0aaef1e94b5b8ee8f1a72e25f5e35e6e5.jpg"),
    7, 4.5, width=5.8)


# =====================================================================
# SLIDE 7: V1 → V2 → V3 Evolution
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, PURPLE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u6a21\u578b\u6f14\u8fdb: V1 \u2192 V2 \u2192 V3",
         font_size=32, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 4, 0.04, PURPLE_ACC)

# V1 column
add_rect(slide, 0.3, 1.3, 3.8, 5.7, RGBColor(0x2a, 0x1a, 0x1a))
add_text(slide, 0.5, 1.4, 3.5, 0.5, "V1: \u5d4c\u5957\u8fd0\u52a8\u94fe",
         font_size=18, color=RED_ACC, bold=True)
add_multiline(slide, 0.5, 1.9, 3.4, 4.8, [
    ("\u67b6\u6784", 14, ORANGE_ACC, True),
    ("plate0: freejoint", 12, LIGHT_GRAY, False),
    ("plates 1-5: \u5d4c\u5957\u5b50\u4f53", 12, LIGHT_GRAY, False),
    ("\u6bcf\u6bb5: slide + 2 hinge", 12, LIGHT_GRAY, False),
    ("8 cable strips/seg", 12, LIGHT_GRAY, False),
    ("", 6, None, False),
    ("\u53c2\u6570", 14, ORANGE_ACC, True),
    ("bend_stiff = 1e8", 12, LIGHT_GRAY, False),
    ("twist_stiff = 4e7", 12, LIGHT_GRAY, False),
    ("muscle_force = 28N", 12, LIGHT_GRAY, False),
    ("yaw_stiff = 50 Nm/rad", 12, RED_ACC, False),
    ("", 6, None, False),
    ("\u95ee\u9898", 14, RED_ACC, True),
    ("\u2718 \u5d4c\u5957\u94fe\u8026\u5408\u5168\u90e8DOF", 12, RED_ACC, False),
    ("\u2718 \u65e0\u6cd5\u72ec\u7acb\u62ac\u8d77\u5355\u6bb5", 12, RED_ACC, False),
    ("\u2718 \u9ad8yaw\u521a\u5ea6\u963b\u6b62\u8f6c\u5f2f", 12, RED_ACC, False),
    ("\u2718 12.4mm/s, 12.2mm\u4fa7\u79fb", 12, RED_ACC, False),
])

# V2 column
add_rect(slide, 4.5, 1.3, 3.8, 5.7, RGBColor(0x1a, 0x1a, 0x2a))
add_text(slide, 4.7, 1.4, 3.5, 0.5, "V2: \u72ec\u7acb\u677f\u4f53",
         font_size=18, color=BLUE_ACC, bold=True)
add_multiline(slide, 4.7, 1.9, 3.4, 4.8, [
    ("\u67b6\u6784", 14, ORANGE_ACC, True),
    ("\u6bcf\u677f\u72ec\u7acb worldbody \u5b50\u4f53", 12, LIGHT_GRAY, False),
    ("\u6bcf\u677f 6 DOF (3\u6ed1\u52a8+3\u94f0\u94fe)", 12, LIGHT_GRAY, False),
    ("\u677f\u95f4\u8f6f\u7126\u63a5 (soft weld)", 12, LIGHT_GRAY, False),
    ("Per-DOF\u521a\u5ea6\u63a7\u5236", 12, GREEN_ACC, False),
    ("", 6, None, False),
    ("\u53c2\u6570", 14, ORANGE_ACC, True),
    ("muscle_force = 38N", 12, LIGHT_GRAY, False),
    ("plate_stiff_x = 500", 12, LIGHT_GRAY, False),
    ("plate_stiff_yaw = 100", 12, RED_ACC, False),
    ("", 6, None, False),
    ("\u6539\u8fdb", 14, GREEN_ACC, True),
    ("\u2714 \u5355\u6bb5\u72ec\u7acb\u62ac\u8d77", 12, GREEN_ACC, False),
    ("\u2714 \u8f74\u5411\u81ea\u7531, \u4fa7\u5411\u7ea6\u675f", 12, GREEN_ACC, False),
    ("\u2714 8.9mm/s, 0.0mm\u4fa7\u79fb", 12, GREEN_ACC, False),
    ("", 6, None, False),
    ("\u6b8b\u7559\u95ee\u9898", 14, RED_ACC, True),
    ("\u2718 yaw_stiff=100 \u4ecd\u6b62\u8f6c\u5f2f", 12, RED_ACC, False),
    ("\u2718 \u65e0\u73af\u808c/\u8f6c\u5411\u808c", 12, RED_ACC, False),
])

# V3 column
add_rect(slide, 8.7, 1.3, 4.3, 5.7, RGBColor(0x1a, 0x2a, 0x1a))
add_text(slide, 8.9, 1.4, 4, 0.5, "V3: \u79bb\u6563\u6b65\u6001+\u5706\u5f27\u8fd0\u52a8",
         font_size=18, color=GREEN_ACC, bold=True)
add_multiline(slide, 8.9, 1.9, 4.0, 4.8, [
    ("\u67b6\u6784", 14, ORANGE_ACC, True),
    ("\u72ec\u7acb\u677f\u4f53 + implicitfast\u79ef\u5206", 12, LIGHT_GRAY, False),
    ("Ring muscles (1/seg)", 12, GREEN_ACC, False),
    ("Diagonal steering (2/seg)", 12, GREEN_ACC, False),
    ("\u7efc\u7ebf\u7ea6\u675f\u53ef\u914d\u7f6e: weld/connect", 12, GREEN_ACC, False),
    ("", 6, None, False),
    ("\u53c2\u6570", 14, ORANGE_ACC, True),
    ("muscle_force = 50N", 12, LIGHT_GRAY, False),
    ("twist_stiff = 2e6 (\u964d20\u00d7)", 12, GREEN_ACC, False),
    ("plate_stiff_yaw = 0 (\u5173\u952e!)", 12, GREEN_ACC, False),
    ("timestep = 0.5ms (\u534a)", 12, LIGHT_GRAY, False),
    ("", 6, None, False),
    ("\u6210\u679c", 14, GREEN_ACC, True),
    ("\u2714 35\u4e2a\u6267\u884c\u5668 (20\u8f74+5\u73af+10\u8f6c)", 12, GREEN_ACC, False),
    ("\u2714 \u5706\u5f27: R=1.3m, aspect=0.96", 12, GREEN_ACC, False),
    ("\u2714 \u7ba1\u9053: 88.7\u00b0\u88ab\u52a8\u8f6c\u5f2f", 12, GREEN_ACC, False),
    ("\u2714 327\u4f53, 1116 DOF", 12, GREEN_ACC, False),
])

# Arrow between columns
add_text(slide, 4.1, 3.5, 0.5, 0.5, "\u2192", font_size=28, color=WHITE, bold=True)
add_text(slide, 8.3, 3.5, 0.5, 0.5, "\u2192", font_size=28, color=WHITE, bold=True)


# =====================================================================
# SLIDE 8: V3 Architecture & Single Segment
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "V3\u6a21\u578b\u67b6\u6784 & \u5355\u6bb5\u9a8c\u8bc1",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 4.5, 0.04, BLUE_ACC)

# Model specs table
spec_data = [
    ["\u53c2\u6570", "\u503c"],
    ["\u6bb5\u6570 Segments", "5 (\u52a0 6 \u677f)"],
    ["\u677f\u76f4\u5f84", "44 mm"],
    ["\u6bb5\u957f", "65 mm"],
    ["\u603b\u4f53\u957f", "~325 mm"],
    ["\u7efc\u7ebf\u63d2\u4ef6", "mujoco.elasticity.cable"],
    ["\u6267\u884c\u5668", "20\u8f74\u5411 + 5\u73af + 10\u8f6c\u5411 = 35"],
    ["DOF", "1116 (\u7efc\u7ebf) / 36 (\u65e0\u7efc)"],
    ["\u4f53\u6570", "327 (\u7efc\u7ebf) / 7 (\u65e0\u7efc)"],
    ["\u65f6\u95f4\u6b65", "0.5 ms (implicitfast)"],
]
add_table(slide, 0.8, 1.3, 5.5, 3.5, 10, 2, spec_data, font_size=12)

# Single segment results
add_text(slide, 0.8, 5.0, 5.5, 0.5,
         "\u5355\u6bb5\u6a21\u5f0f\u9a8c\u8bc1 (Cable Model, 228 DOF)",
         font_size=16, color=BLUE_ACC, bold=True)
cb_data = [
    ["\u6a21\u5f0f", "\u0394Y(mm)", "\u0394X(mm)", "Yaw(\u00b0)", "Ring\u0394R"],
    ["State1 \u6536\u7f29", "-23.66", "0", "0", "+12.55mm"],
    ["State2 \u5de6\u5f2f", "-7.64", "+11.51", "-24.36", "+5.65mm"],
    ["State3 \u53f3\u5f2f", "-7.41", "-11.02", "+23.73", "+5.63mm"],
    ["Ring \u5f84\u5411", "+0.69", "0", "0", "-12.54mm"],
]
add_table(slide, 0.8, 5.4, 6, 1.8, 5, 5, cb_data, font_size=12)

# Right: model images
add_image_safe(slide, os.path.join(IMG_V3, "snap_bend1e8_weld_3q.png"),
               7, 1.3, width=5.8)

# vs Paper comparison
add_rect(slide, 7, 4.5, 5.8, 2.7, DARK_GRAY)
add_text(slide, 7.2, 4.6, 5.5, 0.5, "vs Zhan 2019 \u8bba\u6587\u5bf9\u6bd4",
         font_size=16, color=ORANGE_ACC, bold=True)
add_multiline(slide, 7.2, 5.1, 5.3, 2.0, [
    ("State1: \u6536\u7f29 36.4% (\u8bba\u6587: 42%, \u03b1=0.58)", 13, LIGHT_GRAY, False),
    ("State2: yaw = 24.4\u00b0 (\u8bba\u6587: \u03b8\u2080=18\u00b0)", 13, LIGHT_GRAY, False),
    ("L/R\u5bf9\u79f0\u6027: 0.974 (\u63a5\u8fd1\u5b8c\u7f8e\u5bf9\u79f0)", 13, LIGHT_GRAY, False),
    ("Ring\u6536\u7f29: -52.5% \u534a\u5f84 (\u4f53\u79ef\u5b88\u6052)", 13, LIGHT_GRAY, False),
    ("", 6, None, False),
    ("\u5339\u914d\u826f\u597d \u2014 \u6a21\u578b\u6709\u6548!", 15, GREEN_ACC, True),
])


# =====================================================================
# SLIDE 9: Experiment Results Overview
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u5b9e\u9a8c\u7ed3\u679c\u603b\u89c8  Experiment Results",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 4.5, 0.04, BLUE_ACC)

exp_data = [
    ["\u5b9e\u9a8c", "\u6b65\u6001", "\u7ed3\u679c", "\u72b6\u6001"],
    ["\u76f4\u7ebf\u722c\u884c Rectilinear", "{0,0,2|1}", "~18 mm/s (\u7efc\u7ebf\u6a21\u578b)", "\u2705"],
    ["\u5706\u5f27\u8fd0\u52a8 Circular arc", "{1,0,1|1}", "R=1.3m, aspect=0.96", "\u2705"],
    ["\u5dee\u5206\u8f6c\u5411 Differential", "{1,0,1|1} diff", "6.5\u00b0/s, \u4ec5\u65e0\u7efc\u6a21\u578b", "\u2705"],
    ["\u4f53\u529b\u8f6c\u5411 Body-force", "body_head 0.3N", "R=1.3m, 555s/\u5468", "\u2705"],
    ["\u7ba1\u9053\u722c\u884c (weld)", "{0,0,2|1} pipe", "2.3\u00b0 in 40s", "\u274c"],
    ["\u7ba1\u9053\u722c\u884c (connect 5e7)", "{0,0,2|1} pipe", "69.5\u00b0 in 35s", "\u2705"],
    ["\u7ba1\u9053\u722c\u884c (connect 1e7)", "{0,0,2|1} pipe", "88.7\u00b0 in 40s", "\u2705\u2705"],
]
add_table(slide, 0.5, 1.3, 12.3, 3.3, len(exp_data), 4, exp_data, font_size=13)

# Key discovery boxes
add_rect(slide, 0.5, 4.9, 5.8, 2.3, DARK_GRAY)
add_text(slide, 0.7, 5.0, 5.5, 0.5,
         "\u53d1\u73b0\u2460: plate_stiff_x = 500 \u963b\u6b62\u6240\u6709\u8f6c\u5f2f",
         font_size=15, color=RED_ACC, bold=True)
add_multiline(slide, 0.7, 5.4, 5.4, 1.5, [
    ("\u9ed8\u8ba4\u4fa7\u5411\u521a\u5ea6 500 N/m \u963b\u6b62\u6a2a\u5411\u8fd0\u52a8", 13, LIGHT_GRAY, False),
    ("\u89e3\u51b3: plate_stiff_x = 0 (\u81ea\u7531\u4fa7\u5411)", 13, GREEN_ACC, True),
    ("\u7ecf\u8fc7 5 \u6b21\u5b9e\u9a8c\u624d\u53d1\u73b0\u8fd9\u4e2a\u95ee\u9898!", 13, ORANGE_ACC, False),
])

add_rect(slide, 7, 4.9, 5.8, 2.3, DARK_GRAY)
add_text(slide, 7.2, 5.0, 5.5, 0.5,
         "\u53d1\u73b0\u2461: \u7efc\u7ebf<weld>\u4e0d\u80fd\u8f6c\u5f2f",
         font_size=15, color=RED_ACC, bold=True)
add_multiline(slide, 7.2, 5.4, 5.4, 1.5, [
    ("8\u6761\u94a2\u7efc\u7ebf + <weld> = \u521a\u6027\u7ba1\u4f53", 13, LIGHT_GRAY, False),
    ("bend_stiff=1e8 \u963b\u6b62\u677f\u95f4\u8f6c\u52a8", 13, LIGHT_GRAY, False),
    ("\u89e3\u51b3: <connect> \u7ea6\u675f (Riddle\u7684\u542f\u793a!)", 13, GREEN_ACC, True),
])


# =====================================================================
# SLIDE 10: The Breakthrough — <connect> vs <weld>
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, GREEN_ACC)
add_text(slide, 0.8, 0.4, 12, 0.6,
         "\u7a81\u7834: <connect> vs <weld> \u7ea6\u675f",
         font_size=30, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 5, 0.04, GREEN_ACC)

# Weld side
add_rect(slide, 0.5, 1.4, 5.8, 3.5, RGBColor(0x2a, 0x1a, 0x1a))
add_text(slide, 0.7, 1.5, 5.5, 0.5, "<weld> (6-DOF \u521a\u6027\u7ea6\u675f)",
         font_size=20, color=RED_ACC, bold=True)
add_multiline(slide, 0.7, 2.0, 5.3, 2.8, [
    ("\u9501\u5b9a\u5168\u90e8 6 DOF: 3\u5e73\u79fb + 3\u65cb\u8f6c", 14, LIGHT_GRAY, False),
    ("\u7efc\u7ebf\u7aef\u70b9\u56fa\u5b9a\u4e8e\u677f\u4f53\u65b9\u5411", 14, LIGHT_GRAY, False),
    ("8\u6761\u7efc\u7ebf \u00d7 weld = \u521a\u6027\u7ba1\u72b6\u7ed3\u6784", 14, LIGHT_GRAY, False),
    ("\u4f53\u4e0d\u80fd\u5728\u677f\u8fde\u63a5\u5904\u5f2f\u66f2", 14, LIGHT_GRAY, False),
    ("", 6, None, False),
    ("\u7ba1\u9053\u7ed3\u679c:", 15, RED_ACC, True),
    ("  2.3\u00b0 / 40s (\u9700\u8981 ~30\u5206\u949f\u8fbe 90\u00b0)", 15, RED_ACC, False),
])

# Connect side
add_rect(slide, 7, 1.4, 5.8, 3.5, RGBColor(0x1a, 0x2a, 0x1a))
add_text(slide, 7.2, 1.5, 5.5, 0.5, "<connect> (3-DOF \u7403\u94f0\u5173\u8282)",
         font_size=20, color=GREEN_ACC, bold=True)
add_multiline(slide, 7.2, 2.0, 5.3, 2.8, [
    ("\u4ec5\u7ea6\u675f\u4f4d\u7f6e (3 DOF), \u65cb\u8f6c\u81ea\u7531", 14, LIGHT_GRAY, False),
    ("\u7efc\u7ebf\u7aef\u70b9\u53ef\u81ea\u7531\u65cb\u8f6c", 14, LIGHT_GRAY, False),
    ("\u4f53\u53ef\u5728\u6bcf\u4e2a\u8fde\u63a5\u5904\u88ab\u52a8\u5f2f\u66f2", 14, LIGHT_GRAY, False),
    ("bend_stiff \u964d 2\u00d7 (1e8 \u2192 5e7)", 14, LIGHT_GRAY, False),
    ("", 6, None, False),
    ("\u7ba1\u9053\u7ed3\u679c:", 15, GREEN_ACC, True),
    ("  69.5\u00b0 / 35s (bend=5e7, \u7a33\u5b9a)", 15, GREEN_ACC, False),
    ("  88.7\u00b0 / 40s (bend=1e7, \u7efc\u7ebf\u53d8\u5f62)", 15, GREEN_ACC, False),
])

# Snapshot comparison
add_text(slide, 0.5, 5.1, 6, 0.5,
         "\u89c6\u89c9\u5bf9\u6bd4 (t=15s, \u7ba1\u9053\u5185\u722c\u884c)",
         font_size=16, color=ORANGE_ACC, bold=True)

# Snapshot images
add_image_safe(slide, os.path.join(IMG_V3, "snap_bend1e8_weld_3q.png"),
               0.3, 5.5, width=3.0)
add_text(slide, 0.3, 7.0, 3, 0.3, "weld+1e8: 0\u00b0, \u5b8c\u7f8e\u7efc\u7ebf\u5f62\u72b6",
         font_size=10, color=MID_GRAY)

add_image_safe(slide, os.path.join(IMG_V3, "snap_bend5e7_connect_3q.png"),
               3.5, 5.5, width=3.0)
add_text(slide, 3.5, 7.0, 3, 0.3, "connect+5e7: 9.1\u00b0, \u7efc\u7ebf\u7565\u677e",
         font_size=10, color=MID_GRAY)

add_image_safe(slide, os.path.join(IMG_V3, "snap_bend1e7_connect_3q.png"),
               6.7, 5.5, width=3.0)
add_text(slide, 6.7, 7.0, 3, 0.3, "connect+1e7: 12.8\u00b0, \u7efc\u7ebf\u5d29\u6e83",
         font_size=10, color=MID_GRAY)

# Big comparison text
add_text(slide, 9.8, 5.5, 3.2, 1.5,
         "38\u00d7\n\u63d0\u5347",
         font_size=36, color=GREEN_ACC, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 11: Pipe Crawl Setup & Channel Design
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u7ba1\u9053\u722c\u884c: \u901a\u9053\u8bbe\u8ba1  Channel Design",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 4, 0.04, BLUE_ACC)

ch_data = [
    ["\u53c2\u6570", "\u503c", "\u7406\u7531"],
    ["\u901a\u9053\u5bbd\u5ea6", "56 mm", "44mm\u86f7\u4f53 + 6mm\u95f4\u9699"],
    ["\u5899\u58c1\u9ad8\u5ea6", "55 mm", "\u8986\u76d6\u86f7\u4f53 (max z=45mm)"],
    ["\u76f4\u7ebf\u6bb5\u957f", "400 mm", "~1.2\u500d\u4f53\u957f"],
    ["\u5f2f\u66f2\u534a\u5f84", "200 mm", "\u4e2d\u7b49\u96be\u5ea6"],
    ["\u5f2f\u66f2\u5206\u6bb5", "16\u6bb5", "\u5149\u6ed1\u5f27\u8fd1\u4f3c"],
    ["\u5929\u82b1\u677f\u9ad8\u5ea6", "55 mm", "\u9632\u6b62\u722c\u51fa"],
    ["\u5929\u82b1\u677f\u6469\u64e6", "0.01", "\u8fd1\u65e0\u6469\u64e6"],
    ["\u5899\u58c1\u6469\u64e6", "1.0", "\u6b63\u5e38\u6469\u64e6"],
]
add_table(slide, 0.8, 1.4, 5.5, 3.5, 10, 3, ch_data, font_size=12)

# Collision groups
add_text(slide, 0.8, 5.2, 5.5, 0.5, "\u78b0\u649e\u7ec4\u7b56\u7565",
         font_size=16, color=ORANGE_ACC, bold=True)
cg_data = [
    ["\u5143\u7d20", "contype", "conaffinity", "\u63a5\u89e6\u5bf9\u8c61"],
    ["\u677f\u4f53 Plate", "3", "3", "\u5730\u9762+\u5899\u58c1"],
    ["\u7efc\u7ebf Cable", "4", "4", "\u4ec5\u5730\u9762"],
    ["\u5730\u9762 Ground", "7", "7", "\u5168\u90e8"],
    ["\u5899\u58c1 Wall", "0", "3", "\u677f\u4f53+\u811a"],
]
add_table(slide, 0.8, 5.6, 5.5, 1.8, 5, 4, cg_data, font_size=11)

# Channel diagram
add_rect(slide, 7, 1.3, 5.8, 5.8, RGBColor(0x22, 0x22, 0x35))
add_multiline(slide, 7.3, 1.5, 5.2, 5.0, [
    ("\u901a\u9053\u5e03\u5c40 (\u4fef\u89c6)", 16, BLUE_ACC, True),
    ("", 6, None, False),
    ("    \u2554\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2557   \u2190 outer wall", 14, LIGHT_GRAY, False),
    ("    \u2551  Straight 2   \u2551", 14, LIGHT_GRAY, False),
    ("    \u2551  (\u51fa\u53e3, +X)   \u2551", 14, LIGHT_GRAY, False),
    ("    \u2560\u2550\u2550\u2550\u2557           \u2551   \u2190 inner corner", 14, LIGHT_GRAY, False),
    ("    \u2551   \u2551  90\u00b0 bend \u2551", 14, LIGHT_GRAY, False),
    ("    \u2551   \u255a\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u255d", 14, LIGHT_GRAY, False),
    ("    \u2551       \u2191", 14, LIGHT_GRAY, False),
    ("    \u2551  Straight 1", 14, LIGHT_GRAY, False),
    ("    \u2551  (\u5165\u53e3, +Y)", 14, LIGHT_GRAY, False),
    ("    \u2551       \u2191", 14, LIGHT_GRAY, False),
    ("    \u255a\u2550\u2550\u2550\u2550\u2550\u2550\u2550\u255d", 14, LIGHT_GRAY, False),
    ("   \u86f7\u8d77\u59cb\u4f4d\u7f6e", 14, GREEN_ACC, True),
], font_name="Consolas")


# =====================================================================
# SLIDE 12: Approaches Explored
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u63a2\u7d22\u8fc7\u7a0b: \u6240\u6709\u5c1d\u8bd5\u7684\u65b9\u6848",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 4.5, 0.04, BLUE_ACC)

approach_data = [
    ["#", "\u65b9\u6848", "\u8f6c\u5411\u89d2(\u00b0)", "\u95ee\u9898", "\u72b6\u6001"],
    ["1", "\u65e0\u7efc\u7ebf+\u65e0\u5929\u82b1\u677f", "N/A", "\u86f7\u98de\u51fa\u5899\u58c1 (z=64mm)", "\u274c"],
    ["2", "\u65e0\u7efc\u7ebf+\u5929\u82b1\u677f (f=0.3)", "~5\u00b0", "\u5929\u82b1\u677f\u6469\u64e6\u8fc7\u5927", "\u274c"],
    ["3", "\u65e0\u7efc\u7ebf+\u5929\u82b1\u677f (f=0.01)", "89.3\u00b0", "\u6210\u529f\u4f46\u65e0\u7efc\u7ebf (\u4e0d\u771f\u5b9e)", "\u26a0\ufe0f"],
    ["4", "\u7efc\u7ebf+weld (\u9ed8\u8ba4)", "2.3\u00b0", "8\u00d7weld = \u521a\u6027\u7ba1", "\u274c"],
    ["5", "\u7efc\u7ebf+connect (bend=1e8)", "-0.6\u00b0", "\u4ecd\u7136\u592a\u521a", "\u274c"],
    ["6", "\u7efc\u7ebf+connect (bend=1e6)", "-14mm", "QACC\u4e0d\u7a33\u5b9a, \u540e\u9000", "\u274c"],
    ["7", "\u7efc\u7ebf+connect (bend=1e7)", "88.7\u00b0", "\u7efc\u7ebf\u5d29\u6e83\u4f46\u80fd\u8f6c\u5f2f", "\u26a0\ufe0f"],
    ["8", "\u7efc\u7ebf+connect (bend=5e7)", "51.4\u00b0/30s", "\u6700\u4f73\u89c6\u89c9/\u8f6c\u5f2f\u5e73\u8861", "\u2705"],
    ["9", "\u7efc\u7ebf+connect (bend=1e8, solref=0.02)", "46\u00b0/30s", "\u592a\u6162 (19 sps)", "\u274c"],
]
add_table(slide, 0.3, 1.3, 12.7, 4.5, 11, 5, approach_data, font_size=12)

add_rect(slide, 0.5, 6.0, 12.3, 1.2, DARK_GRAY)
add_multiline(slide, 0.7, 6.1, 12, 1.0, [
    ("\u7ecf\u8fc7 10+ \u6b21\u8fed\u4ee3\u627e\u5230\u6700\u4f73\u65b9\u6848!", 18, ORANGE_ACC, True),
    ("\u5173\u952e\u7ec4\u5408: <connect>\u7ea6\u675f + bend_stiff=5e7 + \u78b0\u649e\u7ec4\u5206\u79bb + \u5c01\u95ed\u901a\u9053 + solref=0.005", 14, GREEN_ACC, False),
])


# =====================================================================
# SLIDE 13: Turning Progress Timeline
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u88ab\u52a8\u8f6c\u5f2f\u65f6\u95f4\u7ebf  Turning Timeline",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 3.5, 0.04, BLUE_ACC)

# Best config (5e7) timeline
add_text(slide, 0.8, 1.3, 5, 0.5,
         "connect + bend=5e7 (solref=0.005)",
         font_size=18, color=GREEN_ACC, bold=True)
tl5e7 = [
    ["\u65f6\u95f4", "\u5934\u90e8\u4f4d\u7f6e (mm)", "\u822a\u5411\u89d2 (\u00b0)", "\u9636\u6bb5"],
    ["t=0s", "(0, 325)", "0.0", "\u76f4\u7ebf\u6bb5"],
    ["t=10s", "(5.9, 460)", "0.9", "\u8fdb\u5165\u5f2f\u9053"],
    ["t=20s", "(90, 587)", "20.6", "\u5f2f\u9053\u4e2d"],
    ["t=30s", "(207, 622)", "51.4", "\u9a76\u51fa\u5f2f\u9053"],
]
add_table(slide, 0.8, 1.8, 5.5, 2.0, 5, 4, tl5e7, font_size=13)

# 1e7 reference
add_text(slide, 0.8, 4.2, 5, 0.5,
         "\u53c2\u8003: connect + bend=1e7 (solref=0.005)",
         font_size=16, color=BLUE_ACC, bold=True)
tl1e7 = [
    ["\u65f6\u95f4", "\u5934\u90e8\u4f4d\u7f6e (mm)", "\u822a\u5411\u89d2 (\u00b0)", "\u5907\u6ce8"],
    ["t=10s", "(20, 510)", "3.4", ""],
    ["t=20s", "(113, 601)", "26.7", ""],
    ["t=30s", "(261, 622)", "64.5", ""],
    ["t=40s", "(506, 608)", "88.7", "\u7efc\u7ebf\u5d29\u6e83\u4f46\u5b8c\u6210\u8f6c\u5f2f"],
]
add_table(slide, 0.8, 4.6, 6, 2.0, 5, 4, tl1e7, font_size=12)

# Right: pipe crawl progression frames (black steel strips, no clipping)
add_image_safe(slide, os.path.join(IMG_V3, "pipe_frame_v3_6_0.png"),
               6.8, 1.2, width=3.1)
add_text(slide, 6.8, 3.55, 3.1, 0.3, "t\u22488s: \u76f4\u7ebf\u6bb5\u8d77\u59cb",
         font_size=10, color=MID_GRAY)
add_image_safe(slide, os.path.join(IMG_V3, "pipe_frame_v3_6_1.png"),
               10.0, 1.2, width=3.1)
add_text(slide, 10.0, 3.55, 3.1, 0.3, "t\u224828s: \u8fdb\u5165\u5f2f\u9053 (25\u00b0)",
         font_size=10, color=MID_GRAY)
add_image_safe(slide, os.path.join(IMG_V3, "pipe_frame_v3_6_2.png"),
               6.8, 4.0, width=3.1)
add_text(slide, 6.8, 6.35, 3.1, 0.3, "t\u224852s: \u5f2f\u9053\u51fa\u53e3 (64\u00b0)",
         font_size=10, color=MID_GRAY)
add_image_safe(slide, os.path.join(IMG_V3, "pipe_frame_v3_6_3.png"),
               10.0, 4.0, width=3.1)
add_text(slide, 10.0, 6.35, 3.1, 0.3, "t\u224872s: \u5b8c\u6210\u8f6c\u5f2f (87\u00b0)",
         font_size=10, color=MID_GRAY)


# =====================================================================
# SLIDE 14: Comparison with Riddle 2025
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u5bf9\u6bd4: \u6211\u4eec\u7684\u6a21\u578b vs Riddle 2025",
         font_size=28, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 5, 0.04, BLUE_ACC)

comp_data = [
    ["", "Riddle 2025", "\u6211\u4eec (V3)"],
    ["\u7269\u7406\u5f15\u64ce", "MuJoCo", "MuJoCo"],
    ["\u7ed3\u6784", "12\u87ba\u65cb PE\u7ba1", "8\u5e73\u884c\u94a2\u7efc\u7ebf"],
    ["\u6750\u6599", "LDPE (0.3 GPa)", "\u5f39\u7c27\u94a2 (~200 GPa)"],
    ["\u9876\u70b9\u5173\u8282", "<connect> \u7403\u94f0", "<connect> (\u5347\u7ea7!)"],
    ["solref", "0.001", "0.005\u20130.01"],
    ["\u4f53\u76f4\u5f84", "320 mm", "44 mm"],
    ["\u6bb5\u6570", "6", "5"],
    ["\u7ba1\u9053\u7c7b\u578b", "\u5706\u5f62\u7ba1 (\u77e9\u5f62box)", "\u5c01\u95ed\u901a\u9053 (\u5899+\u5929\u82b1\u677f)"],
    ["\u6d4b\u8bd5ROC", "0.3\u20130.9 m", "0.2 m"],
    ["\u8f6c\u5f2f\u7ed3\u679c", "~90\u00b0 \u88ab\u52a8", "51.4\u201388.7\u00b0 \u88ab\u52a8"],
    ["\u9a71\u52a8", "Hill\u808c\u8089\u6a21\u578b", "\u8f74\u5411+\u73af+\u8f6c\u5411\u808c\u8089"],
    ["\u6b65\u6001", "3\u00d71\u6ce2 (ramp-hold-ramp)", "{0,0,2|1} \u79bb\u6563\u72b6\u6001\u673a"],
    ["\u4e34\u754cROC", "0.45 m", "\u5f85\u6d4b"],
    ["\u4eff\u771f\u901f\u5ea6", "45\u2013105s/s\u4eff\u771f", "~80 steps/s"],
]
add_table(slide, 0.8, 1.3, 11.5, 5.5, 16, 3, comp_data, font_size=12)

add_text(slide, 0.8, 6.9, 11.5, 0.5,
         "\u5173\u952e\u542f\u793a: <connect> \u7ea6\u675f\u662f MuJoCo \u4e2d\u88ab\u52a8\u8f6c\u5f2f\u7684\u5fc5\u8981\u6761\u4ef6",
         font_size=16, color=GREEN_ACC, bold=True, alignment=PP_ALIGN.CENTER)


# =====================================================================
# SLIDE 15: Next Steps
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, ORANGE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6,
         "\u4e0b\u4e00\u6b65\u8ba1\u5212  Next Steps",
         font_size=32, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 2, 0.04, ORANGE_ACC)

# Short term
add_rect(slide, 0.5, 1.4, 5.8, 2.8, RGBColor(0x1a, 0x2a, 0x1a))
add_text(slide, 0.7, 1.5, 5.5, 0.5, "\u8fd1\u671f  Short-term",
         font_size=22, color=GREEN_ACC, bold=True)
add_bullet_list(slide, 0.7, 2.0, 5.4, 2.0, [
    "\u626b\u63cf\u5f2f\u66f2\u534a\u5f84: [100, 150, 200, 300, 500] mm",
    "\u627e\u5230\u6211\u4eec\u6a21\u578b\u7684\u4e34\u754c ROC",
    "\u626b\u63cf bend_stiff: [1e6, 5e6, 1e7, 5e7, 1e8]",
    "\u4f18\u5316 solref \u53c2\u6570\u4ee5\u63d0\u9ad8\u7a33\u5b9a\u6027",
    "\u6d4b\u91cf\u901f\u5ea6 vs ROC \u5173\u7cfb",
], font_size=14, bullet_color=GREEN_ACC)

# Medium term
add_rect(slide, 7, 1.4, 5.8, 2.8, RGBColor(0x1a, 0x1a, 0x2a))
add_text(slide, 7.2, 1.5, 5.5, 0.5, "\u4e2d\u671f  Medium-term",
         font_size=22, color=BLUE_ACC, bold=True)
add_bullet_list(slide, 7.2, 2.0, 5.4, 2.0, [
    "\u4e3b\u52a8\u8f6c\u5f2f: \u5dee\u5206\u808c\u8089\u63a7\u5236",
    "\u88ab\u52a8+\u4e3b\u52a8\u8f6c\u5411\u7ec4\u5408",
    "\u591a\u79cd\u5f2f\u9053: S\u578b, T\u578b\u4ea4\u53c9\u53e3",
    "\u7ba1\u5185\u969c\u788d\u7269\u89c4\u907f",
    "\u6b65\u6001\u53c2\u6570\u4f18\u5316",
], font_size=14, bullet_color=BLUE_ACC)

# Long term
add_rect(slide, 0.5, 4.5, 12.3, 2.7, RGBColor(0x2a, 0x1a, 0x2a))
add_text(slide, 0.7, 4.6, 5.5, 0.5, "\u8fdc\u671f\u613f\u666f  Long-term",
         font_size=22, color=ORANGE_ACC, bold=True)
add_bullet_list(slide, 0.7, 5.1, 5.4, 2, [
    "\u95ed\u73af\u63a7\u5236: \u63a5\u89e6\u4f20\u611f\u5668\u58c1\u58c1\u68c0\u6d4b",
    "SNS (\u5408\u6210\u795e\u7ecf\u7cfb\u7edf) \u63a7\u5236\u5668",
    "Sim-to-Real \u8f6c\u79fb\u5230\u7269\u7406 CMMWorm",
    "\u53ef\u53d8\u521a\u5ea6\u63a7\u5236",
], font_size=14, bullet_color=ORANGE_ACC)
add_bullet_list(slide, 7.2, 5.1, 5.4, 2, [
    "\u591a\u73af\u5883\u57fa\u51c6: \u5e73\u5730/\u7ba1\u9053/\u788e\u77f3",
    "RL\u6b65\u6001\u4f18\u5316 (PPO/SAC)",
    "\u4e0e Riddle SNS \u63a7\u5236\u5668\u5bf9\u6bd4",
    "\u8bba\u6587: \u88ab\u52a8+\u4e3b\u52a8\u8f6c\u5f2f\u5206\u6790",
], font_size=14, bullet_color=ORANGE_ACC)


# =====================================================================
# SLIDE 16: Summary
# =====================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, 0, 0, 13.333, 0.08, BLUE_ACC)
add_text(slide, 0.8, 0.4, 10, 0.6, "\u603b\u7ed3  Summary",
         font_size=32, color=WHITE, bold=True)
add_rect(slide, 0.8, 1.0, 1.8, 0.04, BLUE_ACC)

add_multiline(slide, 0.8, 1.5, 11.5, 5.5, [
    ("\u2776  \u6784\u5efa 5\u6bb5\u8718\u86f7\u673a\u5668\u4eba MuJoCo \u6a21\u578b, \u7ecf\u5386 V1\u2192V2\u2192V3 \u4e09\u4ee3\u6f14\u8fdb", 17, WHITE, False),
    ("     V1:\u5d4c\u5957\u94fe(\u6f02\u79fb) \u2192 V2:\u72ec\u7acb\u677f(\u7a33\u5b9a) \u2192 V3:\u79bb\u6563\u6b65\u6001+\u8f6c\u5f2f", 13, MID_GRAY, False),
    ("", 6, None, False),
    ("\u2777  \u5355\u6bb5\u9a8c\u8bc1\u4e0e Zhan 2019 \u8bba\u6587\u5339\u914d\u826f\u597d", 17, WHITE, False),
    ("     \u6536\u7f29 36.4% (paper:42%), \u5f2f\u66f2 24\u00b0 (paper:18\u00b0), \u5bf9\u79f0\u6027 0.974", 13, MID_GRAY, False),
    ("", 6, None, False),
    ("\u2778  \u6df1\u5165\u5206\u6790 Riddle 2025: \u53d1\u73b0 <connect> \u7ea6\u675f\u662f\u88ab\u52a8\u8f6c\u5f2f\u7684\u5173\u952e", 17, BLUE_ACC, False),
    ("     LDPE\u7f51\u683c+\u7403\u94f0\u5173\u8282, \u4e34\u754cROC=0.45m, \u5173\u8282\u6469\u64e6\u4f7f\u7b49\u65480.3GPa\u21922GPa", 13, MID_GRAY, False),
    ("", 6, None, False),
    ("\u2779  \u5c06\u7efc\u7ebf-\u677f\u7ea6\u675f\u4ece <weld> \u5347\u7ea7\u4e3a <connect>", 17, GREEN_ACC, True),
    ("     \u7ed3\u5408 bend_stiff \u964d\u4f4e + \u78b0\u649e\u7ec4\u5206\u79bb + \u5c01\u95ed\u901a\u9053", 13, MID_GRAY, False),
    ("", 6, None, False),
    ("\u277a  \u5b9e\u73b0 69.5\u201388.7\u00b0 \u88ab\u52a8\u8f6c\u5f2f (\u4ece 2.3\u00b0) \u2014 30\u00d7 \u63d0\u5347", 20, GREEN_ACC, True),
    ("     \u4ec5\u7528\u76f4\u7ebf\u6b65\u6001, \u65e0\u4e3b\u52a8\u8f6c\u5411, \u7eaf\u9760\u7ba1\u58c1\u7ea6\u675f\u5b9e\u73b0\u65b9\u5411\u6539\u53d8", 13, MID_GRAY, False),
])

# References
add_rect(slide, 0, 6.8, 13.333, 0.7, RGBColor(0x15, 0x15, 0x25))
add_multiline(slide, 0.8, 6.85, 12, 0.6, [
    ("References:  Zhan et al. 2019 IJRR  |  Riddle et al. 2025 Bioinspir. Biomim.  |  "
     "MuJoCo (mujoco.org)  |  GitHub: sriddle97/3D-Soft-Worm-Robot-Model", 11, MID_GRAY, False),
])


# ─── Save ────────────────────────────────────────────────────────────
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"Presentation saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
sz = os.path.getsize(OUT)
print(f"Size: {sz/1024/1024:.2f} MB")
